from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Define slide layouts (0: Title Slide, 1: Title and Content, 6: Blank)
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Helper function to add a slide with title and bullet points
def add_content_slide(title, bullets):
    slide = prs.slides.add_slide(content_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Clear default text
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.LEFT

# Slide 1: Title Slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "Advancements in VLIW DSP Architectures:\nOptimizing Design and Enhancing Security"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)

subtitle = slide.placeholders[1]
subtitle.text = "Paul Sebastien\nTransilvania University of Brasov\n[Conference Name], June 2025"
subtitle.text_frame.paragraphs[0].font.size = Pt(20)
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

# Slide 2: Introduction
add_content_slide(
    "Introduction",
    [
        "VLIW DSPs: Critical for embedded systems (smartphones, IoT, 5G)",
        "Survey synthesizes advancements in design and security",
        "Objective: Optimize performance, power, and security",
        "Why it matters: Balances flexibility and efficiency",
        "Comprehensive analysis for academics and practitioners"
    ]
)

# Slide 3: What Are VLIW DSPs?
add_content_slide(
    "What Are VLIW DSPs?",
    [
        "Leverage instruction-level parallelism (ILP) via static scheduling",
        "Multiple sub-instructions per VLIW (e.g., SWIFT: 4 sub-instructions)",
        "Features: SIMD datapaths, multi-banked memory, simple hardware",
        "Evolution: From TMS320C6x (1990s) to PACDSP, SWIFT",
        "[Placeholder: Diagram of VLIW architecture]"
    ]
)

# Slide 4: Architectural Advancements
add_content_slide(
    "Architectural Advancements",
    [
        "Compressed Instruction Sets: RVC-VOI reduces code size by 25%",
        "Distributed Registers: PACDSP cuts latency by 46.9%",
        "SIMD & Memory: SWIFT’s 4 datapaths for communication",
        "Scalable Designs: PACDSP (0.08 mW/MIPS), SWIFT (SoC)",
        "[Placeholder: Chart comparing code size/latency]"
    ]
)

# Slide 5: Performance Optimization
add_content_slide(
    "Performance Optimization",
    [
        "Compiler: TMS320C6416T -O3 cuts execution by 96.2%",
        "Delay Slot Scheduling: SuperV-DSP boosts speed by 63%",
        "Hand-Coded Assembly: H.263 encoder 61x faster",
        "On-chip SRAM placement: 29x speedup",
        "[Placeholder: Bar graph of performance gains]"
    ]
)

# Slide 6: Power Efficiency
add_content_slide(
    "Power Efficiency",
    [
        "DVFS & Power Gating: PACDSP at 46.60 mW (230 MHz)",
        "Configurable Memory: Saves 25% energy",
        "Synthesis Analysis: 20% energy reduction",
        "Clock Gating: 15% power savings",
        "Trade-off: -O3 increases power by 24.4%",
        "[Placeholder: Power vs. Performance graph]"
    ]
)

# Slide 7: Security Enhancements
add_content_slide(
    "Security Enhancements",
    [
        "Secure Inline Assembly: LLVM-based, prevents code injection",
        "Functional Verification: 40% faster bug detection",
        "Limitations: No defenses for side-channel attacks",
        "Need for secure boot, memory encryption",
        "[Placeholder: Security architecture diagram]"
    ]
)

# Slide 8: Applications
add_content_slide(
    "Applications",
    [
        "Multimedia: TMS320C6700 H.263 encoder (61x speedup)",
        "Communications: SWIFT for IEEE 802.11, 4G",
        "Quantum Simulation: Elbrus reduces latency by 40%",
        "IoT: RVC-VOI’s 25% code reduction for sensors",
        "[Placeholder: Application domain icons]"
    ]
)

# Slide 9: Future Directions
add_content_slide(
    "Future Directions",
    [
        "Security: Hardware-based solutions (secure boot)",
        "Applications: AI, 5G signal processing",
        "Standardization: Platform-agnostic optimizations",
        "Verification: Scalable test generators, simulators",
        "[Placeholder: Roadmap timeline]"
    ]
)

# Slide 10: Conclusion & Q&A
add_content_slide(
    "Conclusion & Q&A",
    [
        "VLIW DSPs: High performance, low power, flexible",
        "Advancements in architecture, optimization, security",
        "Applications: Multimedia, IoT, quantum simulation",
        "Future: Address security, explore AI/5G, standardize",
        "Questions? Thank you!",
        "References: Hsieh et al. [1], Ren et al. [2], Banerjee et al. [3]"
    ]
)

# Save the presentation
prs.save("VLIW_DSP_Conference_Presentation.pptx")